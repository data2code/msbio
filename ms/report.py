#!/usr/bin/env python
from __future__ import absolute_import
from __future__ import print_function
from __future__ import division
import numpy as np
#import mahotas
import scipy.misc
import imageio
from PyPDF2 import PdfFileWriter, PdfFileReader, PdfFileMerger
import os
import tempfile
import pandas as pd
import util
import time
import shutil
import zipfile
import copy
from pptx import Presentation
from pptx.dml.color import RGBColor
import glob
import re
from six.moves import range

def make_icon(s_png, width=100, height=100, s_out=None):
    import scipy.misc
    if width is None and height is None:
        util.error_msg('At least one must not be None: width and height')
    #M=mahotas.imread(s_png)
    M=imageio.imread(s_png)
    h,w,z=M.shape
    if width is None:
        scale=height*1.0/h
    elif height is None:
        scale=width*1.0/w
    else:
        scale_h=height*1.0/h
        scale_w=width*1.0/w
        scale=min(scale_h, scale_w)
    h=int(h*scale+0.5)
    w=int(w*scale+0.5)
    Mr=scipy.misc.imresize(M, (h,w))
    h2,w2,z2=Mr.shape
    M2=np.ones([height, width, 3], dtype='1u1')*255
    offset0=(width-w2)/2
    offset1=(height-h2)/2
    M2[offset1:offset1+h2, offset0:offset0+w2,0]=Mr[:,:,0]
    M2[offset1:offset1+h2, offset0:offset0+w2,1]=Mr[:,:,1]
    M2[offset1:offset1+h2, offset0:offset0+w2,2]=Mr[:,:,2]
    if s_out is None:
        s_out, s_ext=os.path.splitext(s_png)
        s_out+="_thumbnail.png"
    #mahotas.imsave(s_out, M2)
    imageio.imsave(s_out, M2.astype(np.uint8))
# example make_icon('src.png', width=50, height=80, s_out='src_thumbnail.png')

def make_legend_by_cluster(n_cluster, s_out, max_width=300, max_cluster=0, s_label="Cluster"):
    """s_label will be prefix, Cluster 1, Cluster 2...
    if s_label is a list, len(s_label)==n_cluster, we will use that instead"""
    from . import biolists
    m=n_cluster
    if max_cluster>0:
        m=min(n_cluster, max_cluster)
    S_color=biolists.CytoPlot.get_qualitative_colors(m)
    if type(s_label) is list:
        S_label=s_label
        if len(S_label)>m:
            S_label=S_label[:m]
            S_label.append('Other')
            S_color.append('#eeeeee')
    else:
        S_label=['%s %d' % (s_label, i+1) for i in range(m)]
        if m<n_cluster:
            S_color.append('#eeeeee')
            S_label.append('%s >%d' % (s_label, m))
    font_size=14
    return make_legend(S_color, S_label, s_out, max_width=max_width)

def load_font_prop():
    from matplotlib import font_manager as fm
    s_file=os.path.join(os.path.dirname(__file__), "report", "arial.ttf")
    prop = fm.FontProperties(fname=s_file)
    return prop

def load_font():
    from PIL import ImageFont
    font = ImageFont.truetype(os.path.join(os.path.dirname(__file__), "report", "arial.ttf"), 12)
    return font

def make_legend(S_color, S_label, s_out, max_width=300):
    import PIL
    #from PIL import ImageFont
    from PIL import Image
    from PIL import ImageDraw
    font_size=14
    n=len(S_color)
    h=n*font_size # each line has 14 pixels
    w=0
    for x in S_label:
        w=max(w, len(x)*font_size*1.0)
    w+=2
    w=min(int(w), max_width)
    img=Image.new("RGB", (w,h),(255,255,255))
    draw = ImageDraw.Draw(img)
    font = load_font() #ImageFont.truetype(os.path.join(os.path.dirname(__file__), "report", "arial.ttf"), 12)
    for i,s_label in enumerate(S_label):
        clr=S_color[i]
        y=font_size*(i)-1
        draw.rectangle([(1,font_size*i+1),(10,font_size*i+10)], fill=clr)
        draw.text((13, y), s_label,(0,0,0), font=font)
    draw = ImageDraw.Draw(img)
    img.save(s_out)

def make_legend_png(s_png, S_png=[], l_clean_up=True):
    #s_out, s_ext=os.path.splitext(s_png)
    #s_out=s_out+'.png'
    #if os.path.exists(s_out):
    #    os.remove(s_out)
    #M=mahotas.imread(s_png)
    if not os.path.exists(s_png):
        util.warn_msg('File not found in make_legend_png: %s' % s_png)
        return
    try:
        M=imageio.imread(s_png)
        h1,w1,x1 = M.shape
    except:
        util.warn_msg('Cannot load image: %s ' % s_png)
        return
    img_factor=255 if M.dtype==np.uint8 else 1
    h=0
    S_legend=[]
    S_png2=[]
    for x in S_png:
        if not os.path.exists(x):
            util.warn_msg('File not found in make_legend_png: %s' % x)
            continue
        # Add to lower-right corner, bottom up
        X=imageio.imread(x)
        S_legend.append(X)
        S_png2.append(x)
    room=10 # give extra 10px blank space to the left of the legend
    if len(S_legend):
        w=0
        for Q in S_legend:
            w=max(w, Q.shape[1])+room
        if w1<w: # initial image is too small
            M2=np.ones([h1,w,x1])*img_factor
            M2[:,0:w1]=M
            M=M2
            h1,w1,x1 = M.shape
        # find out the extra width required to expand the image, so that legend does not cover up existing pixels
        width=2 # minimum 2 pixel margin
        h=0
        for Q in S_legend:
            h2,w2,z=Q.shape
            w2+=room
            if h1-h<0: break
            height=min(h1-h, h2)
            r=M[max(h1-h-h2,0):h1-h,max(w1-w2,0):w1,0]
            g=M[max(h1-h-h2,0):h1-h,max(w1-w2,0):w1,1]
            b=M[max(h1-h-h2,0):h1-h,max(w1-w2,0):w1,2]
            bg=np.logical_and(r==255, g==255, b==255)
            xb=0
            X=np.sum(bg, axis=0)
            for i in range(r.shape[1]-1,-1,-1):
                if X[i]<height:
                    xb=i
                    break
            width=max(width, xb)
            h+=h2

        # Expand to make sure legend will never overlap the image content
        M2=np.ones([h1,w1+width,x1])*img_factor
        M2[:,0:w1]=M
        h1,w1,x1 = M2.shape
        h=0
        for X in S_legend:
            h2,w2,Z=X.shape
            if x1>Z:
                X=np.dstack([X, np.ones([h2,w2])*img_factor])
            elif x1<Z:
                X=X[:,:,:x1]
            if h1-h<=0: break
            if h1-h-h2<0: # happens when there are too many gene list, legend is too long
                #if h1-h<=0: break
                M2[0:h1-h,w1-w2:]=X[h2-(h1-h):h2,:]
                break
            else:
                M2[h1-h-h2:h1-h,w1-w2:]=X
            h+=h2

        imageio.imsave(s_png, M2.astype(np.uint8))
        if l_clean_up:
            for x in S_png2:
                os.remove(x)

def make_legend_pdf(s_pdf, S_png=[], l_clean_up=True):
    """Add an optional Watermark and legend png (or more), filenames stored in S_png, generate a Watermark.pdf"""
    from reportlab.pdfgen import canvas
    #if s_wm is None:
    #    s_wm=os.path.join(os.path.dirname(__file__), "watermark", "Watermark.pdf")
    if os.path.exists(s_pdf):
        os.remove(s_pdf)
    from reportlab.lib.pagesizes import letter
    c= canvas.Canvas(s_pdf, pagesize=letter)
    w,h=c._pagesize
    w=int(w)
    # place stuff in the center
    w=w/2
    h=h/2
    l_has_page=False
    for x in S_png:
        # Add to lower-right corner, bottom up
        #X=mahotas.imread(x)
        if not os.path.exists(x):
            util.warn_msg('File not found in make_legend_pdf: %s' % x)
            continue
        X=imageio.imread(x)
        h2,w2,Z=X.shape
        c.drawImage(x, w-w2,h)
        h+=h2
        l_has_page=True
        if l_clean_up:
            os.remove(x)
    if l_has_page:
        c.save()
# ://bitbucket.org/rptlab/reportlab

#add_legend_pdf('NewWaterMark.pdf', ['Watermark.png', 'output/test_legend.png'])

def add_watermark_png(s_png, s_out="out", s_wm=None, s_legend=None, l_clean_up=True):
    s_out, s_ext=os.path.splitext(s_out)
    s_out=s_out+'.png'
    if os.path.exists(s_out):
        os.remove(s_out)
    #M=mahotas.imread(s_png)
    if not os.path.exists(s_png):
        util.warn_msg('File not found in make_watermark_png: %s' % s_png)
        return

    M=imageio.imread(s_png)
    if s_wm is None:
        s_wm=os.path.join(os.path.dirname(__file__), "watermark", "Watermark.png")
    #W=mahotas.imread(s_wm)
    if not os.path.exists(s_wm):
        util.warn_msg('File not found in make_watermark_png: %s' % s_wm)
        return
    W=imageio.imread(s_wm)
    if l_clean_up:
        print("Deleting ... "+s_wm)
        os.remove(s_wm)
    h1,w1,x1 = M.shape
    h2,w2,x2 = W.shape
    img_factor=255 if M.dtype==np.uint8 else 1
    if x1>x2: # missing alpha channel
        W=np.dstack([W, np.ones([h2,w2])*img_factor])
    M[-h2:,-w2:]=W
    if s_legend is not None:
        # add legend
        #L=mahotas.imread(s_legend)
        if not os.path.exists(s_legend):
            util.warn_msg('File not found in make_watermark_png: %s' % s_legend)
        else:
            L=imageio.imread(s_legend)
            h3,w3,x3=L.shape
            if x1>x3:
                L=np.dstack([L, np.ones([h3,w3])*img_factor])
            elif x1<x3:
                L=L[:,:,:x1]
            M[-h2-h3:-h2,-w3:]=L
            if l_clean_up:
                print("Deleting ... "+s_legend)
                os.remove(s_legend)
    #mahotas.imsave(s_out, M)
    imageio.imsave(s_out, M.astype(np.uint8))

def add_watermark_pdf(S_pdf, s_out="out", s_wm=None, l_clean_up=True):
    #https://www.binpress.com/tutorial/manipulating-pdfs-with-python/167
    #import shutil
    #shutil.copyfile(s_out, s_out+".v2.pdf")
    s_out, s_ext=os.path.splitext(s_out)
    s_out=s_out+'.pdf'
    #if os.path.exists(s_out):
    #    os.remove(s_out)
    if type(S_pdf) is str:
        S_pdf=[S_pdf]
    if s_wm is None:
        s_wm=os.path.join(os.path.dirname(__file__), "watermark", "Watermark.pdf")
    if not os.path.exists(s_wm):
        util.warn_msg('File not found in make_watermark_pdf: %s' % s_wm)
        return
    f_in=[]
    f_wm=open(s_wm, 'rb')
    wpdf=PdfFileReader(f_wm)
    watermark=wpdf.getPage(0)
    f_in.append(f_wm)
    # there appears to be issue with Python2, we need to close f_wm after it's been used
    out=PdfFileWriter()
    for s_pdf in S_pdf:
        if not os.path.exists(s_pdf):
            util.warn_msg('File not found in make_watermark_pdf: %s' % s_pdf)
            continue
        f=open(s_pdf, 'rb')
        try:
            ipdf=PdfFileReader(f)
        except:
            util.warn_msg('Cannot read PDF: '+s_pdf)
            continue
        for i in range(ipdf.getNumPages()):
            page=ipdf.getPage(i)
            out.addPage(page)
        f_in.append(f)
    out.addPage(watermark)
    # if s_out is one of S_pdf, it may trigger an error msg as s_pdf is still open
    with open(s_out+".tmp", 'wb') as f:
        out.write(f)
    # incase s_out is the same as one of input S_pdf
    for f in f_in:
        f.close()
    if l_clean_up:
        os.remove(s_wm)
        print("Deleting ... "+s_wm)
        for s_pdf in S_pdf:
            if os.path.exists(s_pdf):
                os.remove(s_pdf)
                print("Deleting ... "+s_pdf)
    shutil.move(s_out+".tmp", s_out)

def crop(s_png, margin=10, left=0, right=0, top=0, bottom=0):
    M=imageio.imread(s_png)
    h,w,z=M.shape
    r=M[:,:,0]
    g=M[:,:,1]
    b=M[:,:,2]
    bg=np.logical_and(r==255, g==255, b==255)
    xa=xb=ya=yb=0
    X=np.sum(bg, axis=0)
    #print h,w,X, len(X)
    for i in range(w):
        if X[i]<h:
            xa=i
            break
    for i in range(w-1,-1,-1):
        if X[i]<h:
            xb=i
            break
    X=np.sum(bg, axis=1)
    for i in range(h):
        if X[i]<w:
            ya=i
            break
    for i in range(h-1,-1,-1):
        if X[i]<w:
            yb=i
            break
    if margin>0 or left>0: xa=max(0,xa-max(margin, left))
    if margin>0 or right>0: xb=min(w-1,xb+max(margin, right)+1)
    if margin>0 or top>0: ya=max(0,ya-max(margin,top))
    if margin>0 or bottom>0: yb=min(h-1,yb+max(margin, bottom)+1)
    if xa==0 and xb==w-1 and ya==0 and yb==h-1:
        #print "No cropping"
        return # no cropping
    else:
        #print "Crop: width %d:%d, height: %d:%d" % (ya, yb, xa, xb)
        M=M[ya:yb+1,xa:xb+1,:]
        imageio.imsave(s_png, M.astype(np.uint8))

##https://github.com/scanny/python-pptx/issues/68

def copy_slide(pres, slide_index):
    def _get_blank_slide_layout(pres):
         layout_items_count = [len(layout.placeholders) for layout in pres.slide_layouts]
         min_items = min(layout_items_count)
         blank_layout_id = layout_items_count.index(min_items)
         return pres.slide_layouts[blank_layout_id]

    blank_slide_layout = _get_blank_slide_layout(pres)
    source=pres.slides[slide_index]
    dest = pres.slides.add_slide(blank_slide_layout)

    for shp in source.shapes:
        el = shp.element
        newel = copy.deepcopy(el)
        dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

    for key, value in source.part.rels.items():
        # Make sure we don't copy a notesSlide relation as that won't exist
        if not "notesSlide" in value.reltype:
            dest.part.rels.add_relationship(value.reltype, value._target, value.rId)
    #print "SOURCE>>>>>>>>>>>>>>>>>>>>>>>"
    #info_shapes(source)
    #print "CLONE>>>>>>>>>>>>>>>>>>>>>>>>"
    #info_shapes(dest)
    #print "============================="



def delete_slide(pres, index):
    """Warning: delete should be the last operation, copy after delete seems to lead to corrupted file"""
    xml_slides = pres.slides._sldIdLst  # pylint: disable=W0212
    #slides = list(xml_slides)
    #xml_slides.remove(slides[index])
    del xml_slides[index]

def move_slide(pres, old_index, new_index):
        xml_slides = pres.slides._sldIdLst  # pylint: disable=W0212
        slides = list(xml_slides)
        #xml_slides.remove(slides[old_index])
        slide=slides[old_index]
        del xml_slides[old_index]
        xml_slides.insert(new_index, slide)

def show_slides(pres):
    xml_slides = pres.slides._sldIdLst  # pylint: disable=W0212
    slides = list(xml_slides)
    for i,sld in enumerate(slides):
        print(i, sld)

def info_shape(shp):
    print("\nShape id:", shp.shape_id, " name:", shp.name, " is_placeholder:", shp.is_placeholder, " shape_type:", shp.shape_type, " left:", shp.left, " top:", shp.top, " height:", shp.height, " width:", shp.width, " has_text_frame:", shp.has_text_frame)
    if shp.has_text_frame:
        print(">>>", shp.text)
    print("\n")

def info_shapes(slide):
    for shp in slide.shapes:
        info_shape(shp)

def get_shape(slide, shp_id):
    for shp in slide.shapes:
        if shp.shape_id==shp_id:
            return shp
    util.error_msg('Shape not found: %d', shp_id)

def set_shape_text(slide, shp_id, s_msg):
    shp=get_shape(slide, shp_id)
    if not shp.has_text_frame: util.error_msg('Shape %d is not a text frame' % shp_id)
    shp.text=s_msg

def get_shape_text(slide, shp_id):
    shp=get_shape(slide, shp_id)
    if not shp.has_text_frame: util.error_msg('Shape %d is not a text frame' % shp_id)
    return shp.text

def replace_note(slide, s_msg=None):
    """Only works if note already exist, we just replace"""
    import xml.etree.ElementTree as ET
    import copy
    NameSpace={
        'a':"http://schemas.openxmlformats.org/drawingml/2006/main",
        'r':"http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        'p':"http://schemas.openxmlformats.org/presentationml/2006/main"
    }

    v=self.slide.part_related_by(NameSpace['r']+'/notesSlide')
    r=ET.fromstring(v._blob)
    notes=r.findall('.//p:txBody', NameSpace)
    for note in notes:
        if note.find('.//a:p/a:fld', NameSpace):
            # maybe a slide # or something else
            continue
        for p in note.findall('./a:p', NameSpace):
            note.remove(p)
        msg=ET.fromstring('<a:p xmlns:a="%s"><a:r><a:rPr lang="en-US" dirty="0" smtClean="0"/><a:t></a:t></a:r><a:endParaRPr lang="en-US" dirty="0"/></a:p>' % NameSpace['a'])
        msg.find('.//a:t', NameSpace).text=s_msg
        note.append(msg)
        break
    v.blob=ET.tostring(r)

EMU=914400 #px/inch
VERSION=4

def get_size():
    h=7.5 #  slide height
    if VERSION>=4:
        # 16:9 13 1/3 x 7.5 inch
        w=13.333 # slide width
        left=int(0.7*EMU)
        width=int(w*EMU-2*left)
        height=int((h-1.9)*EMU)
    else:
        # 4:3 10 x 7.5 inch
        #(left, width, height)=(296260, 8551480, 5078313)
        w=10.0 # slide width
        left=int(0.32*EMU)
        width=int(w*EMU-2*left)
        height=int((h-1.9)*EMU)
    return (left, width, height)


def insert_pct(pct_slide, s_img, top=None):
    """My special slide layout, a title and one big picture"""
    (left, width, height)=get_size()
    if top is None:
        top=1347965+25*914400//72
    center_x=left+width//2
    center_y=top+height//2
    ratio0=width*1.0/height
    if not os.path.exists(s_img):
        util.error_msg('File $s not exist!' % s_img)
    try:
        M=imageio.imread(s_img)
        h,w,z=M.shape
    except:
        util.warn_msg("Cannot read image: "+s_img)
        return
    ratio=w*1.0/h
    if ratio>ratio0:
        height=int(width/ratio)
        top=center_y-height//2
    else:
        width=int(height*ratio)
        left=center_x-width//2
    pct_slide.shapes.add_picture(s_img, left, top, width, height)

def insert_table(slide, t, top=None, col_width=None):
    (left, width, height)=get_size()
    if top is None:
        top=1347965+25*914400//72
    rows=len(t)
    cols=len(t.header())
    table = slide.shapes.add_table(rows+1, cols, left, top, width, height).table

    #table.columns[0].width = Inches(2.0)
    #table.columns[1].width = Inches(4.0)

    # write column headings
    S_header=t.header()
    if col_width is not None:
        width_remain=width
        for k,v in col_width.items():
            i=t.col_index(k)
            table.columns[i].width=int(v)
            width_remain-=v
        w=int(max(width_remain/(cols-len(col_width)), 914400//2))
        for i,x in enumerate(S_header):
            if x not in col_width:
                table.columns[i].width=w
    table.rows[0].height = 14*914400//72

    for i in range(cols):
        table.cell(0, i).text = S_header[i]
        for j in range(rows):
            table.cell(j+1, i).text='' if pd.isnull(t.iat[j, i]) else str(t.iat[j, i])
    if '_Color_' in S_header: # used for color legend
        j=util.index('_Color_', S_header)
        for i in range(1,rows+1):
            table.cell(i, j).fill.solid()
            s_hex=t.iat[i-1,j].replace('#', '')
            if re.search(r'^[A-F0-9]{6}$', s_hex):
                table.cell(i, j).fill.fore_color.rgb=RGBColor.from_string(s_hex)
    return table

# to be used in pptx generation
def replace_txt(s_tmpl, s_old, s_new):
    """s_old, s_new can also be lists"""
    if type(s_old) is str:
        s_tmpl=s_tmpl.replace(s_old, str(s_new))
    else:
        for i,x in enumerate(s_old):
            s_tmpl=s_tmpl.replace(x, str(s_new[i]))
    return s_tmpl

# to be used in pptx generation
def add_image(slide, s_img, top=None):
    if s_img is None: return False
    if not os.path.exists(s_img):
        return False
    insert_pct(slide, s_img, top=top)
    return True

def replace_text_keep_style(shape_obj, s_new):
# see https://stackoverflow.com/questions/37924808/python-pptx-power-point-find-and-replace-text-ctrl-h
# we need to do it at run level
    try:
        l_done=False
        for p in shape_obj.text_frame.paragraphs:
            for r in p.runs:
                r.text='' if l_done else s_new
                l_done=True
    except:
        util.warn_msg('text replace failed: %s into %s' % (s_old, s_new))


def pptx_single(s_out, c_txt=None, c_img=None, c_table=None):
    s_out, s_ext=os.path.splitext(s_out)
    s_out=s_out+'.pptx'
    if os.path.exists(s_out):
        os.remove(s_out)
    s_dir=os.path.dirname(__file__)
    prs=Presentation(os.path.join(s_dir, 'report', 'ReportSingle.template.pptx'))
    #shutil.copytree(os.path.join(s_dir, 'ReportSingle'), s_tmp)

    img={} # list of image files
    #    'membership': os.path.join(s_dir,'report/membership.png'),
    #    'GOBargraph': os.path.join(s_dir,'report/GOBargraph.png'),
    #    'GOColorByCluster': os.path.join(s_dir,'report/ColorByCluster.png')
    #    'GOColorByPValue': os.path.join(s_dir,'report/ColorByPValue.png')
    #    'PPIColorByCluster': os.path.join(s_dir,'report/PPIColorByCluster.png')
    #    'GO_MCODE': os.path.join(s_dir,'report/GO_MCODE.csv')
    if c_img is not None:
        img=c_img
    c_table={} if c_table is None else c_table
    # 'GO_MCODE': os.path.join(s_dir,'report/GO_MCODE.csv')
    txt={
        'date': time.strftime("%b %-d, %Y"),
        '#Total': 'X',
        '#Unique': 'Y',
        'SearchTerm': 'none'
    }
    if c_txt is not None:
        txt.update(c_txt)

    sld_cover=prs.slides[0]
    sld_lists=prs.slides[1]
    sld_mem=prs.slides[2]
    sld_bar=prs.slides[3]
    sld_go_clr=prs.slides[4]
    sld_go_p=prs.slides[5]
    sld_ppi_full=prs.slides[6]
    sld_ppi_clr=prs.slides[7]
    sld_ppi_go=prs.slides[8]
    sld_evi_weight=prs.slides[9]

    if VERSION>=4:
        shp=get_shape(sld_cover, 12)
    else:
        shp=get_shape(sld_cover, 7)
    replace_text_keep_style(shp, txt['date'])

    s_msg=get_shape_text(sld_lists, 4)
    s_msg=replace_txt(s_msg, ['#Total', '#Unique'], [txt['#Total'], txt['#Unique']])
    shp=get_shape(sld_lists, 4)
    replace_text_keep_style(shp, s_msg)

    to_delete=[]

    s_msg=get_shape_text(sld_mem, 4)
    s_msg=replace_txt(s_msg, 'SearchTerm', txt['SearchTerm'])
    set_shape_text(sld_mem, 4, s_msg)
    if not add_image(sld_mem, img.get('membership')):
        to_delete.append(2)
    if not add_image(sld_bar, img.get('GOBargraph')):
        to_delete.append(3)
    if not add_image(sld_go_clr, img.get('GOColorByCluster')):
        to_delete.append(4)
    if not add_image(sld_go_p, img.get('GOColorByPValue')):
        to_delete.append(5)
    S_ppi=img.get('PPIColorByCluster',[])
    if type(S_ppi) is not list: S_ppi=[S_ppi]
    # check if _FINAL_ in S_ppi, if yes, we ignore individual file
    S_ppi_final=[x for x in S_ppi if '/_FINAL_' in x]
    if len(S_ppi_final):
        S_ppi=S_ppi_final
    S_mcode=[x for x in S_ppi if '_MCODE_ALL' in x]
    if len(S_ppi)==0:
        to_delete.append(7)
        to_delete.append(6)
    elif len(S_mcode)==0:
        to_delete.append(7)
    for x in S_ppi:
        if '_MCODE_ALL_' in x:
            if not add_image(sld_ppi_clr, x):
                to_delete.append(7)
        else:
            if not add_image(sld_ppi_full, x):
                to_delete.append(7)
                to_delete.append(6)

    if 'GO_MCODE' in c_table:
        if os.path.exists(c_table['GO_MCODE']):
            t=pd.read_csv(c_table['GO_MCODE'])
            insert_table(sld_ppi_go,t, col_width={'Network':914400*1.2})
        else:
            to_delete.append(8)
    else:
        to_delete.append(8)
    if not add_image(sld_evi_weight, img.get('EvidenceWeight')):
        to_delete.append(9)
    to_delete=util.unique(to_delete)
    if len(to_delete):
        to_delete.sort()
        to_delete.reverse()
        for x in to_delete:
            delete_slide(prs, x)
    prs.save(s_out)

def pptx_multiple(s_out, t_lists, c_txt=None, c_img=None, c_table=None, S_color=None):
    """t_lists must be a dataframe of column: "Name, Total, Unique, Color"
    Sequentially listing values for each gene list, color code is Hex
    """
    s_out, s_ext=os.path.splitext(s_out)
    s_out=s_out+'.pptx'
    if os.path.exists(s_out):
        os.remove(s_out)
    s_dir=os.path.dirname(__file__)
    prs=Presentation(os.path.join(s_dir, 'report', 'ReportMultiple.template.pptx'))

    img={}
    #    'membership': os.path.join(s_dir,'report/membership.png'),
    #    'circos': os.path.join(s_dir,'report/CircosOverlapByGene.png'),
    #    'circos_go': os.path.join(s_dir,'report/CircosOverlapByGO.png'),
    #    'GOBargraph': os.path.join(s_dir,'report/HeatmapSelectedGO.png'),
    #    'GOColorByCluster': 'ColorByCluster.png'
    #    'GOColorByPValue': 'ColorByPValue.png'
    #    'GOColorByCounts': 'ColorByCounts.png'
    #    'PPIColorByCluster': ['Brass_MCODE_ALL_PPIColorByCluster.png', 'Karlas_MCODE_ALL_PPIColorByCluster.png', 'Konig_MCODE_ALL_PPIColorByCluster.png', 'MERGE_MCODE_ALL_PPIColorByCluster.png']
    #    'PPIColorByCounts': ['MERGE_MCODE_ALL_PPIColorByCounts.png','MERGE_PPIColorByCounts.png']
    #    'GO_MCODE': os.path.join(s_dir,'report/GO_MCODE.csv')

    if c_img is not None:
        img=c_img
    c_table={} if c_table is None else c_table
    # 'GO_MCODE': os.path.join(s_dir,'report/GO_MCODE_Top3.csv')

    txt={
        'date': time.strftime("%b %-d, %Y"),
        '#Total': 'X',
        '#Unique': 'Y',
        'SearchTerm': 'none'
    }
    if c_txt is not None:
        txt.update(c_txt)

    n_list=len(t_lists)
    from . import biolists
    if '_Color_' not in t_lists.header():
        if S_color is None:
            S_color=biolists.CytoPlot.get_qualitative_colors(n_list)
        t_lists['_Color_']=S_color

    sld_cover=prs.slides[0]
    sld_lists=prs.slides[1]
    sld_cir=prs.slides[2]
    sld_mem=prs.slides[3]
    sld_cir_go=prs.slides[4]
    sld_hm=prs.slides[5]
    sld_go_clr=prs.slides[6]
    sld_go_p=prs.slides[7]
    sld_go_cnt=prs.slides[8]
    sld_ppi_full=prs.slides[9]
    sld_ppi_clr=prs.slides[10]
    sld_ppi_ann=prs.slides[11]
    sld_ppi_full_merge=prs.slides[12]      # _FINAL_COLOR_BY_COUNT
    sld_ppi_full_merge_clr=prs.slides[13]  # _FINAL_COLOR_BY_MCODE
    sld_ppi_clr_merge=prs.slides[14]       # _MCODE_COLOR_BY_CLUSTER
    sld_ppi_cnt_merge=prs.slides[15]       # _MCODE_COLOR_BY_COUNT
    sld_ppi_ann_merge=prs.slides[16]       # ANNOTATION
    sld_evi_weight=prs.slides[17]          # EVIDENCE WEIGHT

    shp=get_shape(sld_cover, 7)
    replace_text_keep_style(shp, txt['date'])

    insert_table(sld_lists, t_lists) #, top=(975360+20*914400/72))

    to_delete=[]
    if not add_image(sld_cir, img.get('circos')):
        to_delete.append(2)
    #s_msg=get_shape_text(sld_mem, 4)
    #s_msg=replace_txt(s_msg, 'SearchTerm', txt['SearchTerm'])
    #set_shape_text(sld_mem, 4, s_msg)

    shp=get_shape(sld_mem, 4)
    replace_text_keep_style(shp, txt['SearchTerm'])

    if not add_image(sld_mem, img.get('membership')):
        to_delete.append(3)
    if not add_image(sld_cir_go, img.get('circos_go')):
        to_delete.append(4)
    if not add_image(sld_hm, img.get('GOBargraph')):
        to_delete.append(5)
    if not add_image(sld_go_clr, img.get('GOColorByCluster')):
        to_delete.append(6)
    if not add_image(sld_go_p, img.get('GOColorByPValue')):
        to_delete.append(7)
    if not add_image(sld_go_cnt, img.get('GOColorByCounts')):
        to_delete.append(8)
    S_ppi=img.get('PPIColorByCluster',[])
    S_merge=img.get('PPIColorByCounts',[])
    if type(S_merge) is not list: S_merge=[S_merge] # _FINAL network
    S_ppi_full=[] # full network for individual gene list
    S_ppi_mcode=[] # MCODE for individual
    S_ppi_merge_mcode=[] # _FINAL_MCODE
    for x in S_ppi:
        if '_MCODE_ALL' in x:
            if '/_FINAL_MCODE_ALL' in x:
                S_ppi_merge_mcode.append(x)
            else:
                S_ppi_mcode.append(x)
                # sometimes full network is too big to plot, it's missing, we will add a dummie entry
                S_ppi_full.append(x.replace('_MCODE_ALL', ''))
        elif '/_FINAL_PPIColorByCluster' not in x:
            S_ppi_full.append(x)
        else:
            S_merge.append(x)
    S_ppi_full=util.unique(S_ppi_full)
    S_ppi_full.sort()
    #print "======================"
    #print S_ppi_full
    #print S_ppi_mcode
    #print S_ppi_merge_mcode
    #print S_merge
    n_new=0
    t_mcode=pd.DataFrame(data=[], columns=['Network','Annotation'])
    if 'GO_MCODE' in c_table:
        if os.path.exists(c_table['GO_MCODE']):
            t_mcode=pd.read_csv(c_table['GO_MCODE'])
    #top=(1347965+20*914400/72)
    if len(S_ppi_full)==0:
        to_delete.extend([9,10,11,12,13,14,15,16])
        idx=17
    else:
        idx=9
        for i in range(len(S_ppi_full)-1):
            copy_slide(prs, 11)
            move_slide(prs, len(prs.slides)-1, 12)
            copy_slide(prs, 10)
            move_slide(prs, len(prs.slides)-1, 12)
            copy_slide(prs, 9)
            move_slide(prs, len(prs.slides)-1, 12)
            n_new+=3
        for x in S_ppi_full:
            s_msg=os.path.basename(x)
            s_name=s_msg.replace('_PPIColorByCluster.png', '')
            s_msg='Gene List: '+s_name
            set_shape_text(prs.slides[idx], 3, s_msg)
            #print ">>>>>1>>>>>>", x
            if not add_image(prs.slides[idx], x): #, top=top):
                to_delete.extend([idx, idx+1, idx+2])
            else:
                s_mcode=x.replace('_PPIColorByCluster', '_MCODE_ALL_PPIColorByCluster')
                #print ">>>>>2>>>>>>", S_ppi_mcode
                if s_mcode in S_ppi_mcode:
                    set_shape_text(prs.slides[idx+1], 3, s_msg)
                    if not add_image(prs.slides[idx+1], s_mcode): #, top=top):
                        to_delete.extend([idx+1])
                else:
                    to_delete.extend([idx+1])
            #print ">>>>3>>>>>", idx, to_delete
            tmp=t_mcode[t_mcode.Network.apply(lambda x: x.startswith(s_name+'_SUB') or x in (s_name, s_name+'_MCODE_ALL'))].copy()
            if len(tmp):
                set_shape_text(prs.slides[idx+2], 3, s_msg)
                insert_table(prs.slides[idx+2], tmp, col_width={'Network':914400*1.2}) #top=top
            else:
                to_delete.append(idx+2)
            #print "BBBBBBBBBBBBBBBBBBBBBBBBBBBBBBBBBBBBB"
            #info_shapes(prs.slides[idx])
            #info_shapes(prs.slides[idx+1])
            #info_shapes(prs.slides[idx+1])
            #print "EEEEEEEEEEEEEEEEEEEEEEEEEEEEEEEEEEEEEEEE"
            idx+=3

        if len(S_merge):
            l_has_mcode=False
            l_has_full=False
            for x in S_merge:
                if '/_FINAL_PPIColorByCounts' in x:
                    if not add_image(sld_ppi_full_merge, x): #, top=top):
                        to_delete.extend([idx,idx+1])
                    else:
                        l_has_full=True
                if '/_FINAL_PPIColorByCluster' in x:
                    if not add_image(sld_ppi_full_merge_clr, x): #, top=top):
                        to_delete.extend([idx,idx+1])
                    else:
                        l_has_full=True
                if '/_FINAL_MCODE_ALL_PPIColorByCounts' in x:
                    if not add_image(sld_ppi_cnt_merge, x): #, top=top):
                        to_delete.extend([idx+3])
                    else:
                        l_has_mcode=True
                    if len(S_ppi_merge_mcode)==0 or not add_image(sld_ppi_clr_merge, S_ppi_merge_mcode[0]): #, top=top):
                        to_delete.extend([idx+2])
                    else:
                        l_has_mcode=True
            if not l_has_full:
                to_delete.extend([idx, idx+1])
            if not l_has_mcode:
                to_delete.extend([idx+2,idx+3])

            tmp=t_mcode[t_mcode.Network.apply(lambda x: x.startswith("_FINAL_") or x=='_FINAL')].copy()
            if len(tmp):
                insert_table(prs.slides[idx+4], tmp, col_width={'Network':914400*1.2}) #top=top
            else:
                to_delete.append(idx+4)
            idx+=5
        else:
            to_delete.extend(list(range(idx,idx+5)))
    if not add_image(sld_evi_weight, img.get('EvidenceWeight')):
        to_delete.append(idx)

    to_delete=util.unique(to_delete)
    if len(to_delete):
        to_delete=util.unique(to_delete)
        to_delete.sort()
        to_delete.reverse()
        for x in to_delete:
            delete_slide(prs, x)
    prs.save(s_out)

def make_analysis_report(s_src='.', s_out='AnalysisReport.zip'):

    def add_file_to_zip(x, s_file, s_arch):
        if os.path.exists(s_file):
            x.write(s_file, s_arch)

    def add_folder_to_zip(x, s_dir, s_arch_dir):
        if not os.path.exists(s_dir): return
        for root, folder, files in os.walk(s_dir):
            for filename in files:
                if filename.startswith("."): continue
                s_src=os.path.join(root, filename)
                s_dest=os.path.join(s_arch_dir, s_src.replace(s_dir+"/", ''))
                add_file_to_zip(x, s_src, s_dest)

    ar=zipfile.ZipFile(s_out, 'w', zipfile.ZIP_DEFLATED)
    for x in ["README.txt", 'icon/SVG48.png','icon/PDF48.png','icon/CYS48.png','icon/WEB_CYS48.png']:
        ar.write(os.path.join(os.path.dirname(__file__), "report", x), x)
    for x in ['AnalysisReport.pptx','AnalysisReport.html','COVID_map.xlsx']:
        add_file_to_zip(ar, os.path.join(s_src, x), x)
    for x in ['membership.png','membership.pdf']:
        add_file_to_zip(ar, os.path.join(s_src, 'Membership_PieChart', x), os.path.join('Membership_PieChart', x))
    for x in ['Evidence.csv', 'EvidenceWeight.csv', 'EvidenceWeight.png', 'EvidenceWeight.pdf']:
        add_file_to_zip(ar, os.path.join(s_src, 'Evidence', x), os.path.join('Evidence', x))
    for x in ['GO_AllLists.csv','ColorByCluster.png', 'ColorByCluster.pdf', 'ColorByPValue.png', 'ColorByPValue.pdf', 'ColorByCounts.png', 'ColorByCounts.pdf', 'GONetwork.cys', 'GONetwork.xgmml', '_FINAL_GO.csv', 'GeneGo_membership.csv', 'GO_membership.csv', 'GO_membership_parent.csv', 'GONetwork.html', 'GONetwork.js', 'GONetwork.style.js']:
        add_file_to_zip(ar, os.path.join(s_src, 'Enrichment_GO', x), os.path.join('Enrichment_GO', x))
    add_folder_to_zip(ar, os.path.join(s_src, 'Enrichment_QC'), 'Enrichment_QC')
    for x in ['CircosOverlapByGene.png', 'CircosOverlapByGene.svg', 'CircosOverlapByGO.png', 'CircosOverlapByGO.svg']:
        add_file_to_zip(ar, os.path.join(s_src, 'Overlap_circos', x), os.path.join('Overlap_circos', x))
    for suffix in ["", "Top100", "Parent"]:
        for ext in ["cdt","atr","gtr","jtv","png","pdf","csv"]:
            s=os.path.join('Enrichment_heatmap', 'HeatmapSelectedGO'+suffix+"."+ext)
            add_file_to_zip(ar, os.path.join(s_src, s), s)
    for x in ['GO_MCODE_Top3.csv', 'MCODE.csv', '_FINAL_MCODE.csv', 'PPINetwork.html', 'PPINetwork.js', 'PPINetwork.style.js']:
        add_file_to_zip(ar, x, os.path.join('Enrichment_PPI', x))
    for x in glob.glob(s_src+'/Enrichment_PPI/*') + glob.glob(s_src+'/Enrichment_PPI/xgmml/*'):
        if not os.path.splitext(x)[1] in ['.json','.pickle']:
            add_file_to_zip(ar, x, os.path.join('Enrichment_PPI',
                                             ('xgmml/' if '/xgmml/' in x else '') +os.path.basename(x)))
    for x in glob.glob(s_src+'/Enrichment_PieChart/*'):
        add_file_to_zip(ar, x, os.path.join('Enrichment_PieChart', os.path.basename(x)))
    if os.path.exists(os.path.join(s_src, 'Enrichment_GO/GONetwork.html')) or os.path.exists(os.path.join(s_src, 'Enrichment_PPI/PPINetwork.html')):
        import ms.mssetting
        cyjs_src=ms.mssetting.cytoscapejs['SOURCE']
        add_folder_to_zip(ar, cyjs_src, "CyJS")
    ar.close()


if __name__=="__main__":
    t_lists=pd.DataFrame(data={'Name':['Brass','Karlas', 'Konig'], '#Total':[300,200, 158], '#Unique':[300, 198, 158]})
    pptx_multiple(
        'MyReport2.pptx',
        t_lists,
        c_txt={'SearchTerm': 'Invasion'},
        c_img={'membership': 'output/membership.png', 'circos': 'output/CircosOverlapByGene.png', 'circos_go': 'report/CircosOverlapByGO.png', 'GOBargraph': 'output/HeatmapSelectedGO.png', 'GOColorByCluster':'output/ColorByCluster.png', 'GOColorByPValue':'output/ColorByPValue.png', 'GOColorByCounts': 'output/ColorByCounts.png', 'PPIColorByCluster':['output/PPI/Brass_MCODE_ALL_PPIColorByCluster.png', 'output/PPI/Karlas_MCODE_ALL_PPIColorByCluster.png', 'output/PPI/Konig_MCODE_ALL_PPIColorByCluster.png', 'output/PPI/Brass_PPIColorByCluster.png', 'output/PPI/Karlas_PPIColorByCluster.png', 'output/PPI/Konig_PPIColorByCluster.png', 'output/PPI/_FINAL_MCODE_ALL_PPIColorByCluster.png','output/PPI/_FINAL_PPIColorByCluster.png'], 'PPIColorByCounts': ['output/PPI/_FINAL_MCODE_ALL_PPIColorByCounts.png','output/PPI/_FINAL_PPIColorByCounts.png']},
        c_table={'GO_MCODE': 'output/GO_MCODE_Top3.csv'}
    )
    pptx_single(
        'MyReport.pptx',
        c_txt={'#Total': 250, '#Unique':250, 'SearchTerm': 'Invasion'},
        c_img={'membership': 'output/membership.png', 'GOBargraph': 'output/HeatmapSelectedGO.png', 'GOColorByCluster':'output/ColorByCluster.png', 'GOColorByPValue':'output/ColorByPValue.png', 'PPIColorByCluster':['output/PPI/Gene_MCODE_ALL_PPIColorByCluster.png','output/PPI/Gene_PPIColorByCluster.png']},
        c_table={'GO_MCODE': 'output/GO_MCODE_Top3.csv'}
    )
    exit()

    import ms.report as report

    import ms.biolists as bl
    n=10
    S_color=bl.CytoPlot.get_qualitative_colors(n)
    S_label=['Activator List','Below Normal','Control List', 'Delicate List', 'Elegant List', 'Fancy List', 'Glossy List', 'High Signal','Intermediate','Just too long to display list']

    make_legend(S_color, S_label, s_out='output/test_legend.png')

    #S_color=brewer2mpl.get_map("Set1", 'qualitative', 3).hex_colors[:]
    report.add_watermark_png('output/ColorByCounts.png', s_out='output/ColorByCounts_v2.png')
    report.add_watermark_pdf(['output/ColorByCounts.pdf','output/ColorByCluster.pdf','output/ColorByPValue.pdf'], s_out='output/GONetwork.pdf')
